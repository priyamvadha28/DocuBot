# Importing necessary modules
import streamlit as st
from PyPDF2 import PdfReader
import docx
import pptx
import pandas as pd
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import AzureOpenAIEmbeddings
from langchain.chains.question_answering import load_qa_chain
from langchain_openai import AzureChatOpenAI
from langchain.prompts import PromptTemplate

# Initializing the chat model with required configuration details
Chat_Model = {
    "openai_api_type": "azure",  # Specifies that we're using Azure OpenAI
    "deployment_name": "",       # Fill in with your deployment name
    "model_name": "",            # Fill in with the specific model name you're using
    "openai_api_base": "",       # Fill in with your Azure OpenAI base endpoint
    "openai_api_version": "",    # Fill in with the API version you're using
    "openai_api_key": "",        # Fill in with your OpenAI API key
    "temperature": 0             # Temperature setting for the LLM (0 = deterministic)
}

# Define the LLM (Language Model) using AzureChatOpenAI with the above configurations
llm = AzureChatOpenAI(openai_api_type=Chat_Model['openai_api_type'],
                      deployment_name=Chat_Model['deployment_name'],
                      model_name=Chat_Model['model_name'],
                      azure_endpoint=Chat_Model['openai_api_base'],
                      openai_api_version=Chat_Model['openai_api_version'],
                      openai_api_key=Chat_Model['openai_api_key'],
                      temperature=0)

# Setting the header for the Streamlit app
st.header("DocuBot")

# Initialize chat history in session state if it doesn't exist
if "chats" not in st.session_state:
    st.session_state["chats"] = {"history": []}

# Sidebar for uploading documents and viewing chat history
with st.sidebar:
    st.title("Upload document here")
    # File uploader to upload various document types
    file = st.file_uploader("Upload a document and start asking questions", type=["pdf", "docx", "pptx", "xlsx"])

    # Display chat history in an expandable format
    st.subheader("Chat History")
    for message in st.session_state["chats"]["history"]:
        with st.expander(f"{message['role']}"):
            st.write(message['content'])

# Input field for user questions
user_question = st.text_input("Type your question here")

# Initialize vector_store as None, will be used to store document embeddings
vector_store = None

# Function to extract text from the uploaded document based on its type
def extract_text_from_file(file):
    if file.name.endswith('.pdf'):
        pdf_reader = PdfReader(file)
        text = ""
        # Extract text from each page of the PDF
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    elif file.name.endswith('.docx'):
        doc = docx.Document(file)
        # Extract text from all paragraphs in the Word document
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.name.endswith('.pptx'):
        presentation = pptx.Presentation(file)
        text = ""
        # Extract text from all slides in the PowerPoint presentation
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.name.endswith('.xlsx'):
        df = pd.read_excel(file)
        # Convert the Excel data to a string
        return df.to_string(index=False)
    else:
        return ""

# If a file is uploaded, extract text and create vector store from document chunks
if file is not None:
    text = extract_text_from_file(file)

    if text:
        # Split the extracted text into manageable chunks using RecursiveCharacterTextSplitter
        text_splitter = RecursiveCharacterTextSplitter(
            separators="\n",
            chunk_size=1000,
            chunk_overlap=150,
            length_function=len
        )

        chunks = text_splitter.split_text(text)

        # Initialize embeddings using AzureOpenAIEmbeddings
        embeddings = AzureOpenAIEmbeddings(
            openai_api_version="2024-02-01",
            api_key="a3efe3e1fe1f4b6ca6f9d6d1a23f8991",
            azure_endpoint="https://nextgen-azure-openai.openai.azure.com/",
            model="nextgen-embedding-model"
        )
        
        # Create a vector store from the text chunks using FAISS
        vector_store = FAISS.from_texts(chunks, embeddings)
    else:
        st.error("Failed to extract text from the uploaded file.")

# Function to manage the chat session
def chat_session():
    chat_data = st.session_state["chats"]

    if user_question:
        # Store user's question in session state
        chat_data["history"].append({"role": "You", "content": user_question})

        if vector_store:
            # Search for relevant text in the document using the vector store
            match = vector_store.similarity_search(user_question)

            # Prompt template for generating responses from the model
            prompt_template = """
            Here is some relevant information:
            {context}

            Based on this information, answer the following question.
            If you don't know the answer or cannot find it in the document, give the message:
            "The document does not have the necessary information required to answer this question."
            
            If no question is asked, only an empty space is given, give the message:
            "No question is given, please enter a valid question to obtain a reply."
            {question}
            """
            prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])

            # Load the question-answering chain with the defined prompt and LLM
            chain = load_qa_chain(llm, chain_type="stuff", prompt=prompt)
            try:
                # Generate a response using the QA chain
                response = chain.run(input_documents=match, question=user_question)
                # Store bot's response in session state
                chat_data["history"].append({"role": "Bot", "content": response})
                
                # Display only the bot's response in the main section
                st.write(f"**Bot:** {response}")
            except Exception as e:
                # Handle and display errors
                error_message = f"An error occurred: {str(e)}"
                st.error(error_message)
                chat_data["history"].append({"role": "Bot", "content": error_message})
                st.write(f"**Bot:** {error_message}")

# Run the chat session function
chat_session()

# This project is developed and maintained by Priyamvadha Pradeep.